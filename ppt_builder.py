"""
Build a PPTX deck from a list of Excel rows.

For v1 there is one fixed layout (the "Excellent Worklets" template) and one
default Excel-to-PPT mapping. The mapping is data-driven so future versions
can let users supply their own template + mapping config.
"""

import io
from typing import Any, Dict, List, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REQUIRED_COLUMNS = [
    "Worklet Id",
    "Worklet Title",
    "College",
    "Mentor 1",
    "Mentor 2",
    "Brief Description",
    "Initial Deliverables vs Final Output",
    "Q.4 Key Achievements",
    "Team",
    "Mentor Final Remarks",
]


# A mapping value is either a single column name or a list of column names
# whose values are concatenated with a space.
ColumnSpec = Union[str, List[str]]

DEFAULT_MAPPING: Dict[str, ColumnSpec] = {
    "header_title": ["Worklet Id", "Worklet Title"],
    "college": "College",
    "team": "Team",
    "mentor1": "Mentor 1",
    "mentor2": "Mentor 2",
    "problem_statement": "Brief Description",
    "expectation": "Initial Deliverables vs Final Output",
    "achievements": "Q.4 Key Achievements",
    "mentors_remarks": "Mentor Final Remarks",
}


# Colors picked to match the reference template.
COLOR_HEADER_BAR = RGBColor(0xB7, 0xE0, 0xEC)
COLOR_PS_HEADER = RGBColor(0x2F, 0x55, 0x97)
COLOR_EXP_HEADER = RGBColor(0x7F, 0x7F, 0x7F)
COLOR_ACH_HEADER = RGBColor(0x54, 0x82, 0x35)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_LINK_BLUE = RGBColor(0x00, 0x70, 0xC0)


def _solid_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _set_line(shape, rgb: RGBColor, width_pt: float = 0.75) -> None:
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def _hide_line(shape) -> None:
    shape.line.fill.background()


def _write_lines(
    tf,
    lines: List[str],
    *,
    size: int = 11,
    bold: bool = False,
    color: RGBColor = COLOR_BLACK,
    align=PP_ALIGN.LEFT,
    font_name: str = "Calibri",
) -> None:
    tf.word_wrap = True
    tf.clear()
    if not lines:
        lines = [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = color


def _write_text(tf, text: str, **kwargs) -> None:
    _write_lines(tf, str(text).split("\n") if text else [""], **kwargs)


def _add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    return box, tf


def _add_rounded_rect(slide, left, top, width, height, *, fill=None, line_color=COLOR_BLACK, line_width_pt=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.06
    if fill is None:
        shape.fill.background()
    else:
        _solid_fill(shape, fill)
    if line_color is None:
        _hide_line(shape)
    else:
        _set_line(shape, line_color, line_width_pt)
    return shape


def _add_rect(slide, left, top, width, height, *, fill, hide_line: bool = True):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _solid_fill(shape, fill)
    if hide_line:
        _hide_line(shape)
    return shape


def _resolve(row: Dict[str, Any], spec: ColumnSpec) -> str:
    if isinstance(spec, list):
        return " ".join(str(row.get(c, "")).strip() for c in spec if str(row.get(c, "")).strip())
    return str(row.get(spec, "")).strip()


EXCLUDED_MENTOR_NAMES = {"samsung prism"}


def _build_slide(
    prs: Presentation,
    row: Dict[str, Any],
    mapping: Dict[str, ColumnSpec],
    top_title: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header_title = _resolve(row, mapping["header_title"])
    college = _resolve(row, mapping["college"])
    team = _resolve(row, mapping["team"])
    mentor1 = _resolve(row, mapping["mentor1"])
    mentor2 = _resolve(row, mapping["mentor2"])
    problem = _resolve(row, mapping["problem_statement"])
    expectation = _resolve(row, mapping["expectation"])
    achievements = _resolve(row, mapping["achievements"])
    remarks = _resolve(row, mapping["mentors_remarks"])

    mentors = [
        m for m in (mentor1, mentor2)
        if m and m.strip().lower() not in EXCLUDED_MENTOR_NAMES
    ]

    # Top title (configurable per generation)
    _, tf = _add_textbox(slide, Inches(0.2), Inches(0.1), Inches(13.0), Inches(0.55))
    _write_text(tf, top_title, size=24, bold=True)

    # Light-blue header bar
    bar_left, bar_top = Inches(0.2), Inches(0.85)
    bar_width, bar_height = Inches(13.0), Inches(0.95)
    _add_rect(slide, bar_left, bar_top, bar_width, bar_height, fill=COLOR_HEADER_BAR)

    # Worklet Id + Title (left side of bar)
    _, tf = _add_textbox(slide, Inches(0.3), Inches(0.95), Inches(7.5), Inches(0.75))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _write_text(tf, header_title, size=18, bold=True)

    # College + Team (middle of bar)
    _, tf = _add_textbox(slide, Inches(7.9), Inches(0.9), Inches(2.7), Inches(0.9))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _write_lines(
        tf,
        [f"College : {college}", f"Team- {team}"],
        size=14,
        bold=True,
    )

    # Mentors (right side of bar, right-aligned, "Samsung PRISM" filtered out)
    _, tf = _add_textbox(slide, Inches(10.6), Inches(0.9), Inches(2.6), Inches(0.9))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_right = Inches(0.1)
    _write_lines(tf, mentors or [""], size=12, bold=True, align=PP_ALIGN.RIGHT)

    # Three top boxes
    box_top = Inches(1.95)
    box_height = Inches(3.4)
    header_height = Inches(0.4)
    box_w = Inches(4.25)
    lefts = [Inches(0.2), Inches(4.55), Inches(8.9)]
    sections = [
        ("Problem Statement", COLOR_PS_HEADER, problem),
        ("Expectation",       COLOR_EXP_HEADER, expectation),
        ("Achievements",      COLOR_ACH_HEADER, achievements),
    ]
    for left, (label, header_color, body) in zip(lefts, sections):
        _add_rounded_rect(slide, left, box_top, box_w, box_height, fill=None, line_color=COLOR_BLACK)
        hdr = _add_rect(slide, left, box_top, box_w, header_height, fill=header_color)
        hdr.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _write_text(hdr.text_frame, label, size=16, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

        _, body_tf = _add_textbox(
            slide,
            left + Inches(0.1),
            box_top + header_height + Inches(0.05),
            box_w - Inches(0.2),
            box_height - header_height - Inches(0.1),
        )
        _write_text(body_tf, body, size=12)

    # Bottom row
    bot_top = Inches(5.5)
    bot_height = Inches(1.85)

    # Output box (left, static placeholders)
    out_left, out_w = Inches(0.2), Inches(2.7)
    _add_rounded_rect(slide, out_left, bot_top, out_w, bot_height, fill=None, line_color=COLOR_BLACK)
    _, tf = _add_textbox(slide, out_left, bot_top + Inches(0.05), out_w, Inches(0.4))
    _write_text(tf, "Output", size=16, bold=True, align=PP_ALIGN.CENTER)

    img_w, img_h = Inches(0.7), Inches(0.85)
    img_y = bot_top + Inches(0.45)
    for cx in (out_left + Inches(0.35), out_left + Inches(1.55)):
        ph = _add_rect(slide, cx, img_y, img_w, img_h, fill=COLOR_WHITE, hide_line=False)
        _set_line(ph, COLOR_BLACK, 0.5)

    _, tf = _add_textbox(slide, out_left + Inches(0.05), img_y + img_h + Inches(0.02), Inches(1.3), Inches(0.25))
    _write_text(tf, "Commercialization", size=10, align=PP_ALIGN.CENTER)
    _, tf = _add_textbox(slide, out_left + Inches(1.25), img_y + img_h + Inches(0.02), Inches(1.3), Inches(0.25))
    _write_text(tf, "Publications", size=10, align=PP_ALIGN.CENTER)
    _, tf = _add_textbox(slide, out_left + Inches(0.05), img_y + img_h + Inches(0.27), Inches(1.3), Inches(0.25))
    _write_text(tf, "<Commercialization>", size=10, color=COLOR_LINK_BLUE, align=PP_ALIGN.CENTER)
    _, tf = _add_textbox(slide, out_left + Inches(1.25), img_y + img_h + Inches(0.27), Inches(1.3), Inches(0.25))
    _write_text(tf, "<Forum>", size=10, align=PP_ALIGN.CENTER)

    # Results box (matches Output / Mentors Remarks styling, body left as placeholder)
    res_left, res_w = Inches(3.0), Inches(5.7)
    _add_rounded_rect(slide, res_left, bot_top, res_w, bot_height, fill=None, line_color=COLOR_BLACK)
    _, tf = _add_textbox(slide, res_left, bot_top + Inches(0.05), res_w, Inches(0.4))
    _write_text(tf, "Results", size=16, bold=True, align=PP_ALIGN.CENTER)

    # Mentors Remarks box (right, data-driven from "Mentor Final Remarks")
    rem_left, rem_w = Inches(8.9), Inches(4.3)
    _add_rounded_rect(slide, rem_left, bot_top, rem_w, bot_height, fill=None, line_color=COLOR_BLACK)
    _, tf = _add_textbox(slide, rem_left, bot_top + Inches(0.05), rem_w, Inches(0.4))
    _write_text(tf, "Mentors Remarks", size=16, bold=True, align=PP_ALIGN.CENTER)
    _, body_tf = _add_textbox(
        slide,
        rem_left + Inches(0.1),
        bot_top + Inches(0.5),
        rem_w - Inches(0.2),
        bot_height - Inches(0.55),
    )
    _write_text(body_tf, remarks, size=12, bold=True)


DEFAULT_TOP_TITLE = "Excellent Worklets – Feb26-April26"


def build_presentation(
    rows: List[Dict[str, Any]],
    mapping: Dict[str, ColumnSpec] = None,
    top_title: str = DEFAULT_TOP_TITLE,
) -> bytes:
    if mapping is None:
        mapping = DEFAULT_MAPPING

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for row in rows:
        _build_slide(prs, row, mapping, top_title)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
